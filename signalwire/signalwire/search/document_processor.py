"""
Copyright (c) 2025 SignalWire

This file is part of the SignalWire AI Agents SDK.

Licensed under the MIT License.
See LICENSE file in the project root for full license information.
"""

import re
import hashlib
import json
from typing import List, Dict, Any, Optional
from pathlib import Path

# Document processing imports
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    import markdown
except ImportError:
    markdown = None

try:
    from striprtf.striprtf import rtf_to_text
except ImportError:
    rtf_to_text = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from nltk.tokenize import sent_tokenize
    import nltk
    # Ensure NLTK data is available
    try:
        nltk.data.find('tokenizers/punkt')
    except LookupError:
        nltk.download('punkt', quiet=True)
except ImportError:
    sent_tokenize = None
    nltk = None

try:
    import magic
except ImportError:
    magic = None

from .query_processor import preprocess_document_content
from signalwire.core.logging_config import get_logger

logger = get_logger(__name__)

class DocumentProcessor:
    """Enhanced document processor with smart chunking capabilities"""
    
    def __init__(
        self,
        chunking_strategy: str = 'sentence',
        max_sentences_per_chunk: int = 5,
        chunk_size: int = 50,
        chunk_overlap: int = 10,
        split_newlines: Optional[int] = None,
        index_nlp_backend: str = 'nltk',
        verbose: bool = False,
        semantic_threshold: float = 0.5,
        topic_threshold: float = 0.3
    ):
        """
        Initialize document processor

        Args:
            chunking_strategy: Strategy for chunking documents:
                - 'sentence': Sentence-based chunking with overlap
                - 'sliding': Sliding window with word-based chunks
                - 'paragraph': Natural paragraph boundaries
                - 'page': Page-based chunking (for PDFs)
                - 'semantic': Semantic similarity-based chunking
                - 'topic': Topic modeling-based chunking
                - 'qa': Question-answer optimized chunking
                - 'json': JSON structure-aware chunking
                - 'markdown': Markdown structure-aware chunking with code block detection
            max_sentences_per_chunk: For sentence strategy (default: 5)
            chunk_size: For sliding strategy - words per chunk (default: 50)
            chunk_overlap: For sliding strategy - overlap in words (default: 10)
            split_newlines: For sentence strategy - split on multiple newlines (optional)
            index_nlp_backend: NLP backend for indexing (default: 'nltk')
            verbose: Whether to enable verbose logging (default: False)
            semantic_threshold: Similarity threshold for semantic chunking (default: 0.5)
            topic_threshold: Similarity threshold for topic chunking (default: 0.3)
        """
        self.chunking_strategy = chunking_strategy
        self.max_sentences_per_chunk = max_sentences_per_chunk
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.split_newlines = split_newlines
        self.semantic_threshold = semantic_threshold
        self.topic_threshold = topic_threshold
        
        # Legacy support for old character-based chunking
        self.chunk_overlap = chunk_overlap
    
    def create_chunks(self, content: str, filename: str, 
                     file_type: str) -> List[Dict[str, Any]]:
        """
        Create chunks from document content using specified chunking strategy
        
        Args:
            content: Document content (string) - should be the actual content, not a file path
            filename: Name of the file (for metadata)
            file_type: File extension/type
            
        Returns:
            List of chunk dictionaries
        """
        
        # Apply chunking strategy
        if self.chunking_strategy == 'sentence':
            return self._chunk_by_sentences(content, filename, file_type)
        elif self.chunking_strategy == 'sliding':
            return self._chunk_by_sliding_window(content, filename, file_type)
        elif self.chunking_strategy == 'paragraph':
            return self._chunk_by_paragraphs(content, filename, file_type)
        elif self.chunking_strategy == 'page':
            return self._chunk_by_pages(content, filename, file_type)
        elif self.chunking_strategy == 'semantic':
            return self._chunk_by_semantic(content, filename, file_type)
        elif self.chunking_strategy == 'topic':
            return self._chunk_by_topics(content, filename, file_type)
        elif self.chunking_strategy == 'qa':
            return self._chunk_by_qa_optimization(content, filename, file_type)
        elif self.chunking_strategy == 'json':
            return self._chunk_from_json(content, filename, file_type)
        elif self.chunking_strategy == 'markdown':
            # Use markdown-aware chunking for better structure preservation
            return self._chunk_markdown_enhanced(content, filename)
        else:
            # Fallback to sentence-based chunking
            return self._chunk_by_sentences(content, filename, file_type)
    
    def _extract_text_from_file(self, file_path: str) -> Any:
        """Extract text from various file formats"""
        if not magic:
            # Fallback to extension-based detection
            file_path_obj = Path(file_path)
            extension = file_path_obj.suffix.lower()
            
            if extension == '.pdf':
                file_type = 'application/pdf'
            elif extension == '.docx':
                file_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            elif extension in ['.txt', '.md']:
                file_type = 'text/plain'
            elif extension == '.html':
                file_type = 'text/html'
            elif extension == '.rtf':
                file_type = 'application/rtf'
            elif extension == '.xlsx':
                file_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            elif extension == '.pptx':
                file_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            else:
                file_type = 'text/plain'
        else:
            try:
                mime = magic.Magic(mime=True)
                file_type = mime.from_file(file_path)
            except (FileNotFoundError, IOError, OSError):
                # Fall back to extension-based detection if magic can't read the file
                file_path_obj = Path(file_path)
                extension = file_path_obj.suffix.lower()
                ext_map = {
                    '.pdf': 'application/pdf',
                    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    '.txt': 'text/plain', '.md': 'text/plain',
                    '.html': 'text/html',
                    '.rtf': 'application/rtf',
                    '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                }
                file_type = ext_map.get(extension, 'text/plain')

        if 'pdf' in file_type:
            return self._extract_pdf(file_path)
        elif 'vnd.openxmlformats-officedocument.wordprocessingml.document' in file_type:
            return self._extract_docx(file_path)
        elif 'plain' in file_type or 'text' in file_type:
            return self._extract_text(file_path)
        elif 'html' in file_type:
            return self._extract_html(file_path)
        elif 'markdown' in file_type or file_path.endswith('.md'):
            return self._extract_markdown(file_path)
        elif 'rtf' in file_type:
            return self._extract_rtf(file_path)
        elif 'vnd.openxmlformats-officedocument.spreadsheetml.sheet' in file_type:
            return self._extract_excel(file_path)
        elif 'vnd.openxmlformats-officedocument.presentationml.presentation' in file_type:
            return self._extract_powerpoint(file_path)
        else:
            return json.dumps({"error": f"Unsupported file type: {file_type}"})
    
    def _extract_pdf(self, file_path: str):
        """Extract text from PDF files"""
        if not pdfplumber:
            return json.dumps({"error": "pdfplumber not available for PDF processing"})
        
        try:
            with pdfplumber.open(file_path) as pdf:
                pages = []
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        # Remove page number from the beginning
                        text = re.sub(r'^\d+\.\s*', '', text.strip())
                        pages.append(text)
                return pages
        except Exception as e:
            return json.dumps({"error": f"Error processing PDF: {e}"})
    
    def _extract_docx(self, file_path: str):
        """Extract text from DOCX files"""
        if not DocxDocument:
            return json.dumps({"error": "python-docx not available for DOCX processing"})
        
        try:
            doc = DocxDocument(file_path)
            return [para.text for para in doc.paragraphs if para.text.strip()]
        except Exception as e:
            return json.dumps({"error": f"Error processing DOCX: {e}"})
    
    def _extract_text(self, file_path: str):
        """Extract text from plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            return json.dumps({"error": f"Error processing TXT: {e}"})
    
    def _extract_html(self, file_path: str):
        """Extract text from HTML files"""
        if not BeautifulSoup:
            return json.dumps({"error": "beautifulsoup4 not available for HTML processing"})
        
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                soup = BeautifulSoup(file, 'html.parser')
                return soup.get_text(separator='\n')
        except Exception as e:
            return json.dumps({"error": f"Error processing HTML: {e}"})
    
    def _extract_markdown(self, file_path: str):
        """Extract text from Markdown files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
                if markdown and BeautifulSoup:
                    html = markdown.markdown(content)
                    soup = BeautifulSoup(html, 'html.parser')
                    return soup.get_text(separator='\n')
                else:
                    # Fallback to raw markdown
                    return content
        except Exception as e:
            return json.dumps({"error": f"Error processing Markdown: {e}"})
    
    def _extract_rtf(self, file_path: str):
        """Extract text from RTF files"""
        if not rtf_to_text:
            return json.dumps({"error": "striprtf not available for RTF processing"})
        
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return rtf_to_text(file.read())
        except Exception as e:
            return json.dumps({"error": f"Error processing RTF: {e}"})
    
    def _extract_excel(self, file_path: str):
        """Extract text from Excel files"""
        if not load_workbook:
            return json.dumps({"error": "openpyxl not available for Excel processing"})
        
        try:
            wb = load_workbook(file_path)
            sheets_text = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' '.join([str(cell) for cell in row if cell is not None])
                    sheets_text.append(row_text)
            return "\n".join(sheets_text)
        except Exception as e:
            return json.dumps({"error": f"Error processing Excel: {e}"})
    
    def _extract_powerpoint(self, file_path: str):
        """Extract text from PowerPoint files"""
        if not Presentation:
            return json.dumps({"error": "python-pptx not available for PowerPoint processing"})
        
        try:
            prs = Presentation(file_path)
            slides_text = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                slides_text.append("\n".join(slide_text))
            return slides_text
        except Exception as e:
            return json.dumps({"error": f"Error processing PowerPoint: {e}"})
    
    def _chunk_document_aware(self, content: Any, filename: str, file_type: str) -> List[Dict[str, Any]]:
        """Smart chunking for documents with natural structure"""
        chunks = []
        
        if isinstance(content, list):
            # Handle page-based or paragraph-based content (PDF, DOCX, PPTX)
            for i, page_content in enumerate(content):
                if not page_content or not page_content.strip():
                    continue
                
                # For each page/slide, use sentence-based chunking if it's large
                if len(page_content) > self.chunk_size * 6:  # chunk_size is in words, ~6 chars/word
                    page_chunks = self._sentence_based_chunking(
                        page_content, 
                        max_sentences_per_chunk=self._calculate_sentences_per_chunk(page_content)
                    )
                    for j, chunk_content in enumerate(page_chunks):
                        chunks.append(self._create_chunk(
                            content=chunk_content,
                            filename=filename,
                            section=f"Page {i+1}" if file_type == 'pdf' else f"Slide {i+1}" if file_type == 'pptx' else f"Section {i+1}",
                            metadata={'page_number': i+1, 'chunk_index': j}
                        ))
                else:
                    # Small page/slide - keep as single chunk
                    chunks.append(self._create_chunk(
                        content=page_content,
                        filename=filename,
                        section=f"Page {i+1}" if file_type == 'pdf' else f"Slide {i+1}" if file_type == 'pptx' else f"Section {i+1}",
                        metadata={'page_number': i+1}
                    ))
        else:
            # Single text content - use paragraph-aware chunking
            chunks = self._chunk_text_enhanced(content, filename)
        
        return chunks
    
    def _chunk_markdown_enhanced(self, content: str, filename: str) -> List[Dict[str, Any]]:
        """Enhanced markdown chunking with code block detection and rich metadata

        Features:
        - Tracks header hierarchy for section paths
        - Detects code blocks and extracts language
        - Adds 'code' tags to chunks containing code
        - Preserves markdown structure for better search
        """
        chunks = []
        lines = content.split('\n')

        current_section = None
        current_hierarchy = []  # Track header hierarchy
        current_chunk = []
        current_size = 0
        line_start = 1
        in_code_block = False
        code_languages = []  # Track languages in current chunk
        has_code = False

        for line_num, line in enumerate(lines, 1):
            # Check for code block fences
            code_fence_match = re.match(r'^```(\w+)?', line)
            if code_fence_match:
                in_code_block = not in_code_block
                if in_code_block:
                    # Starting code block
                    has_code = True
                    lang = code_fence_match.group(1)
                    if lang and lang not in code_languages:
                        code_languages.append(lang)

            # Check for headers with hierarchy tracking
            header_match = re.match(r'^(#{1,6})\s+(.+)', line) if not in_code_block else None
            if header_match:
                header_level = len(header_match.group(1))
                header_text = header_match.group(2).strip()

                # Save current chunk if it exists
                if current_chunk:
                    chunk_metadata = self._build_markdown_metadata(
                        current_hierarchy, code_languages, has_code
                    )
                    chunks.append(self._create_chunk(
                        content='\n'.join(current_chunk),
                        filename=filename,
                        section=self._build_section_path(current_hierarchy),
                        start_line=line_start,
                        end_line=line_num - 1,
                        metadata=chunk_metadata
                    ))

                # Update hierarchy
                current_hierarchy = current_hierarchy[:header_level-1] + [header_text]
                current_section = header_text
                current_chunk = [line]
                current_size = len(line)
                line_start = line_num
                code_languages = []
                has_code = False

            else:
                current_chunk.append(line)
                current_size += len(line) + 1

                # Check if chunk is getting too large - use smart splitting
                # But don't split inside code blocks
                if current_size >= self.chunk_size * 6 and not in_code_block:  # chunk_size is in words, ~6 chars/word
                    # Try to split at paragraph boundary first
                    split_point = self._find_best_split_point(current_chunk)

                    chunk_to_save = current_chunk[:split_point]
                    chunk_metadata = self._build_markdown_metadata(
                        current_hierarchy, code_languages, has_code
                    )
                    chunks.append(self._create_chunk(
                        content='\n'.join(chunk_to_save),
                        filename=filename,
                        section=self._build_section_path(current_hierarchy),
                        start_line=line_start,
                        end_line=line_start + split_point - 1,
                        metadata=chunk_metadata
                    ))

                    # Start new chunk with overlap
                    overlap_lines = self._get_overlap_lines(chunk_to_save)
                    remaining_lines = current_chunk[split_point:]
                    current_chunk = overlap_lines + remaining_lines
                    current_size = sum(len(line) + 1 for line in current_chunk)
                    line_start = line_start + split_point - len(overlap_lines)
                    # Reset code tracking for new chunk
                    code_languages = []
                    has_code = False

        # Add final chunk
        if current_chunk:
            chunk_metadata = self._build_markdown_metadata(
                current_hierarchy, code_languages, has_code
            )
            chunks.append(self._create_chunk(
                content='\n'.join(current_chunk),
                filename=filename,
                section=self._build_section_path(current_hierarchy),
                start_line=line_start,
                end_line=len(lines),
                metadata=chunk_metadata
            ))

        return chunks
    
    def _chunk_python_enhanced(self, content: str, filename: str) -> List[Dict[str, Any]]:
        """Enhanced Python code chunking with better function/class detection"""
        chunks = []
        lines = content.split('\n')
        
        current_function = None
        current_class = None
        current_chunk = []
        current_size = 0
        line_start = 1
        indent_level = 0
        
        for line_num, line in enumerate(lines, 1):
            # Detect class definitions
            class_match = re.match(r'^(\s*)(class\s+([^(:\s]+))', line)
            if class_match:
                indent = len(class_match.group(1))
                class_name = class_match.group(3)
                
                # Save current chunk if switching context
                if current_chunk and (indent <= indent_level or current_class != class_name):
                    chunks.append(self._create_chunk(
                        content='\n'.join(current_chunk),
                        filename=filename,
                        section=self._build_python_section(current_class, current_function),
                        start_line=line_start,
                        end_line=line_num - 1
                    ))
                    current_chunk = []
                    line_start = line_num
                
                current_class = class_name
                current_function = None
                indent_level = indent
            
            # Detect function definitions
            func_match = re.match(r'^(\s*)(def\s+([^(:\s]+)|async\s+def\s+([^(:\s]+))', line)
            if func_match:
                indent = len(func_match.group(1))
                func_name = func_match.group(3) or func_match.group(4)
                
                # Save current chunk if switching to new function at same or lower level
                if current_chunk and indent <= indent_level:
                    chunks.append(self._create_chunk(
                        content='\n'.join(current_chunk),
                        filename=filename,
                        section=self._build_python_section(current_class, current_function),
                        start_line=line_start,
                        end_line=line_num - 1
                    ))
                    current_chunk = []
                    line_start = line_num
                
                if indent >= indent_level:  # Method within class or nested function
                    current_function = func_name
                else:  # Top-level function
                    current_function = func_name
                    current_class = None
                
                indent_level = indent
            
            current_chunk.append(line)
            current_size += len(line) + 1
            
            # Handle oversized chunks
            if current_size >= self.chunk_size * 6:  # chunk_size is in words, ~6 chars/word
                chunks.append(self._create_chunk(
                    content='\n'.join(current_chunk),
                    filename=filename,
                    section=self._build_python_section(current_class, current_function),
                    start_line=line_start,
                    end_line=line_num
                ))
                
                # Start new chunk with minimal overlap for code
                overlap_lines = current_chunk[-2:] if len(current_chunk) > 2 else current_chunk
                current_chunk = overlap_lines
                current_size = sum(len(line) + 1 for line in overlap_lines)
                line_start = line_num - len(overlap_lines) + 1
        
        # Add final chunk
        if current_chunk:
            chunks.append(self._create_chunk(
                content='\n'.join(current_chunk),
                filename=filename,
                section=self._build_python_section(current_class, current_function),
                start_line=line_start,
                end_line=len(lines)
            ))
        
        return chunks
    
    def _chunk_text_enhanced(self, content: str, filename: str) -> List[Dict[str, Any]]:
        """Enhanced text chunking using sentence-based approach"""
        if isinstance(content, list):
            content = '\n'.join(content)
        
        # Use sentence-based chunking for better coherence
        max_sentences = self._calculate_sentences_per_chunk(content)
        sentences = self._sentence_based_chunking(content, max_sentences)
        
        chunks = []
        for i, chunk_content in enumerate(sentences):
            chunks.append(self._create_chunk(
                content=chunk_content,
                filename=filename,
                section=f"Section {i+1}",
                metadata={'chunk_method': 'sentence_based', 'chunk_index': i}
            ))
        
        return chunks
    
    def _sentence_based_chunking(self, text: str, max_sentences_per_chunk: int, split_newlines: int = 2) -> List[str]:
        """Sentence-based chunking with enhancements"""
        if not sent_tokenize:
            # Fallback to simple splitting
            sentences = text.split('. ')
            sentences = [s.strip() + '.' for s in sentences if s.strip()]
        else:
            sentences = []

            if split_newlines > 0:
                # Create regex pattern for specified number of newlines
                newline_pattern = r'(\n{%d,})' % split_newlines
                parts = re.split(newline_pattern, text)
                
                for part in parts:
                    part = part.strip()
                    if part and not re.match(newline_pattern, part):
                        sentences.extend(sent_tokenize(part))
                    elif re.match(newline_pattern, part):
                        sentences.append(part)
            else:
                sentences = [sentence.strip() for sentence in sent_tokenize(text) if sentence.strip()]

        # Create chunks of sentences with overlap
        chunks = []
        overlap_sentences = max(1, max_sentences_per_chunk // 4)  # 25% overlap
        
        for i in range(0, len(sentences), max_sentences_per_chunk - overlap_sentences):
            chunk_sentences = sentences[i:i + max_sentences_per_chunk]
            if chunk_sentences:
                chunks.append(' '.join(chunk_sentences))
        
        return chunks
    
    def _calculate_sentences_per_chunk(self, text: str) -> int:
        """Calculate optimal sentences per chunk based on average sentence length"""
        if not sent_tokenize:
            # Fallback calculation
            sentences = text.split('. ')
        else:
            sentences = sent_tokenize(text)
            
        if not sentences:
            return 1
        
        avg_sentence_length = sum(len(s) for s in sentences) / len(sentences)
        # Target chunk size divided by average sentence length
        optimal_sentences = max(1, int(self.chunk_size / avg_sentence_length))
        return min(optimal_sentences, 10)  # Cap at 10 sentences for readability
    
    def _build_section_path(self, hierarchy: List[str]) -> str:
        """Build hierarchical section path from header hierarchy"""
        return ' > '.join(hierarchy) if hierarchy else None

    def _build_markdown_metadata(self, hierarchy: List[str], code_languages: List[str], has_code: bool) -> Dict[str, Any]:
        """Build rich metadata for markdown chunks

        Args:
            hierarchy: Current header hierarchy (e.g., ['Installation', 'Requirements', 'Python'])
            code_languages: List of code block languages found in chunk (e.g., ['python', 'bash'])
            has_code: Whether chunk contains any code blocks

        Returns:
            Dictionary with markdown-specific metadata including tags
        """
        metadata = {
            'chunk_type': 'markdown',
        }

        # Add header level metadata
        if hierarchy:
            for i, header in enumerate(hierarchy, 1):
                metadata[f'h{i}'] = header

        # Add code-related metadata
        if has_code:
            metadata['has_code'] = True
            if code_languages:
                metadata['code_languages'] = code_languages

        # Build tags for enhanced searching
        tags = []
        if has_code:
            tags.append('code')
            # Add language-specific tags
            for lang in code_languages:
                tags.append(f'code:{lang}')

        # Add tags for header levels (searchable by section depth)
        if len(hierarchy) > 0:
            tags.append(f'depth:{len(hierarchy)}')

        if tags:
            metadata['tags'] = tags

        return metadata
    
    def _build_python_section(self, class_name: Optional[str], function_name: Optional[str]) -> str:
        """Build section name for Python code"""
        if class_name and function_name:
            return f"{class_name}.{function_name}"
        elif class_name:
            return class_name
        elif function_name:
            return function_name
        else:
            return None
    
    def _find_best_split_point(self, lines: List[str]) -> int:
        """Find the best point to split a chunk (prefer paragraph boundaries)"""
        # Look for empty lines (paragraph boundaries) in the last 25% of the chunk
        start_search = max(1, len(lines) * 3 // 4)
        
        for i in range(len(lines) - 1, start_search - 1, -1):
            if not lines[i].strip():  # Empty line
                return i
        
        # If no paragraph boundary found, split at 75% of chunk size
        return max(1, len(lines) * 3 // 4)
    
    def _create_chunk(self, content: str, filename: str, 
                     section: Optional[str] = None,
                     start_line: Optional[int] = None,
                     end_line: Optional[int] = None,
                     metadata: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        """Create chunk dictionary with enhanced metadata"""
        base_metadata = {
            'file_type': Path(filename).suffix.lstrip('.'),
            'chunk_size': len(content),
            'word_count': len(content.split()),
        }
        
        # Handle sentence count with fallback
        try:
            if sent_tokenize and content.strip():
                base_metadata['sentence_count'] = len(sent_tokenize(content))
            else:
                # Fallback: count sentences by periods
                base_metadata['sentence_count'] = len([s for s in content.split('.') if s.strip()])
        except Exception as e:
            logger.warning(f"Error counting sentences: {e}")
            # Simple fallback: count periods
            base_metadata['sentence_count'] = len([s for s in content.split('.') if s.strip()])
        
        if metadata:
            base_metadata.update(metadata)
        
        return {
            'content': content.strip(),
            'filename': filename,
            'section': section,
            'start_line': start_line,
            'end_line': end_line,
            'metadata': base_metadata
        }
    
    def _get_overlap_lines(self, lines: List[str]) -> List[str]:
        """Get overlap lines for chunk continuity"""
        if not lines:
            return []
        
        # Calculate overlap size in characters
        overlap_chars = self.chunk_overlap
        overlap_lines = []
        char_count = 0
        
        # Take lines from the end until we reach overlap size
        for line in reversed(lines):
            if char_count + len(line) <= overlap_chars:
                overlap_lines.insert(0, line)
                char_count += len(line) + 1
            else:
                break
        
        return overlap_lines
    
    def _chunk_by_sentences(self, content: str, filename: str, file_type: str) -> List[Dict[str, Any]]:
        """Chunk content by sentences with specified max sentences per chunk"""
        if isinstance(content, list):
            content = '\n'.join(content)
        
        # Use sentence-based chunking
        split_newlines = self.split_newlines if self.split_newlines is not None else 2
        sentences = self._sentence_based_chunking(content, self.max_sentences_per_chunk, split_newlines)
        
        chunks = []
        for i, chunk_content in enumerate(sentences):
            chunks.append(self._create_chunk(
                content=chunk_content,
                filename=filename,
                section=f"Section {i+1}",
                metadata={
                    'chunk_method': 'sentence_based', 
                    'chunk_index': i,
                    'max_sentences_per_chunk': self.max_sentences_per_chunk,
                    'split_newlines': split_newlines
                }
            ))
        
        return chunks
    
    def _chunk_by_sliding_window(self, content: str, filename: str, file_type: str) -> List[Dict[str, Any]]:
        """Chunk content using sliding window approach with word-based chunks"""
        if isinstance(content, list):
            content = '\n'.join(content)
        
        # Split content into words
        words = content.split()
        
        if not words:
            return []
        
        chunks = []
        chunk_index = 0
        
        # Create overlapping chunks
        for i in range(0, len(words), self.chunk_size - self.chunk_overlap):
            chunk_words = words[i:i + self.chunk_size]
            if chunk_words:
                chunk_content = ' '.join(chunk_words)
                chunks.append(self._create_chunk(
                    content=chunk_content,
                    filename=filename,
                    section=f"Chunk {chunk_index + 1}",
                    metadata={
                        'chunk_method': 'sliding_window',
                        'chunk_index': chunk_index,
                        'chunk_size_words': self.chunk_size,
                        'overlap_size_words': self.chunk_overlap,
                        'start_word': i,
                        'end_word': i + len(chunk_words)
                    }
                ))
                chunk_index += 1
        
        return chunks
    
    def _chunk_by_paragraphs(self, content: str, filename: str, file_type: str) -> List[Dict[str, Any]]:
        """Chunk content by paragraphs (split on double newlines)"""
        if isinstance(content, list):
            content = '\n'.join(content)
        
        # Split on double newlines to get paragraphs
        paragraphs = re.split(r'\n\s*\n', content)
        paragraphs = [p.strip() for p in paragraphs if p.strip()]
        
        chunks = []
        for i, paragraph in enumerate(paragraphs):
            if paragraph:
                chunks.append(self._create_chunk(
                    content=paragraph,
                    filename=filename,
                    section=f"Paragraph {i+1}",
                    metadata={
                        'chunk_method': 'paragraph_based',
                        'chunk_index': i,
                        'paragraph_number': i + 1
                    }
                ))
        
        return chunks
    
    def _chunk_by_pages(self, content: str, filename: str, file_type: str) -> List[Dict[str, Any]]:
        """Chunk content by pages (for documents that have page boundaries)"""
        if isinstance(content, list):
            # If content is already a list (e.g., from PDF extraction), treat each item as a page
            pages = [str(page).strip() for page in content if str(page).strip()]
        else:
            # For text content, try to detect page boundaries
            # Look for form feed characters or page break indicators
            if '\f' in content:
                pages = content.split('\f')
            elif '---PAGE---' in content:
                pages = content.split('---PAGE---')
            elif re.search(r'\n\s*Page\s+\d+\s*\n', content):
                # Split on "Page N" patterns
                pages = re.split(r'\n\s*Page\s+\d+\s*\n', content)
            else:
                # Fallback: split into roughly equal chunks
                words = content.split()
                words_per_page = max(500, len(words) // 10)  # Aim for ~10 pages
                pages = []
                for i in range(0, len(words), words_per_page):
                    page_words = words[i:i + words_per_page]
                    if page_words:
                        pages.append(' '.join(page_words))
        
        pages = [p.strip() for p in pages if p.strip()]
        
        chunks = []
        for i, page_content in enumerate(pages):
            if page_content:
                chunks.append(self._create_chunk(
                    content=page_content,
                    filename=filename,
                    section=f"Page {i+1}",
                    metadata={
                        'chunk_method': 'page_based',
                        'chunk_index': i,
                        'page_number': i + 1
                    }
                ))
        
        return chunks
    
    def _chunk_by_semantic(self, content: str, filename: str, file_type: str) -> List[Dict[str, Any]]:
        """Chunk based on semantic similarity between sentences"""
        if isinstance(content, list):
            content = '\n'.join(content)
        
        # Get sentences
        if sent_tokenize:
            sentences = sent_tokenize(content)
        else:
            sentences = content.split('. ')
            sentences = [s.strip() + '.' for s in sentences if s.strip()]
        
        if len(sentences) <= 1:
            return [self._create_chunk(content, filename, "Section 1", 
                                     metadata={'chunk_method': 'semantic', 'chunk_index': 0})]
        
        # Generate embeddings for sentences (using the same model as the index)
        try:
            from sentence_transformers import SentenceTransformer
            from sklearn.metrics.pairwise import cosine_similarity
            import numpy as np
            
            from signalwire.search.query_processor import _get_cached_model
            model = _get_cached_model('sentence-transformers/all-mpnet-base-v2')
            embeddings = model.encode(sentences, show_progress_bar=False)
            
            # Calculate similarity between adjacent sentences
            similarities = []
            for i in range(len(embeddings) - 1):
                sim = cosine_similarity([embeddings[i]], [embeddings[i + 1]])[0][0]
                similarities.append(sim)
            
            # Find split points where similarity drops below threshold
            split_points = [0]
            for i, sim in enumerate(similarities):
                if sim < self.semantic_threshold:
                    split_points.append(i + 1)
            split_points.append(len(sentences))
            
            # Create chunks
            chunks = []
            for i in range(len(split_points) - 1):
                start_idx = split_points[i]
                end_idx = split_points[i + 1]
                chunk_sentences = sentences[start_idx:end_idx]
                
                # Ensure minimum chunk size
                if len(chunk_sentences) < 2 and i > 0:
                    # Merge with previous chunk
                    chunks[-1]['content'] += ' ' + ' '.join(chunk_sentences)
                    continue
                    
                chunk_content = ' '.join(chunk_sentences)
                chunks.append(self._create_chunk(
                    content=chunk_content,
                    filename=filename,
                    section=f"Semantic Section {i+1}",
                    metadata={
                        'chunk_method': 'semantic',
                        'chunk_index': i,
                        'semantic_threshold': self.semantic_threshold,
                        'sentence_count': len(chunk_sentences)
                    }
                ))
            
            return chunks if chunks else [self._create_chunk(content, filename, "Section 1",
                                                           metadata={'chunk_method': 'semantic', 'chunk_index': 0})]
            
        except ImportError:
            # Fallback to sentence-based chunking
            return self._chunk_by_sentences(content, filename, file_type)

    def _chunk_by_topics(self, content: str, filename: str, file_type: str) -> List[Dict[str, Any]]:
        """Chunk based on topic changes using keyword analysis"""
        if isinstance(content, list):
            content = '\n'.join(content)
        
        if sent_tokenize:
            sentences = sent_tokenize(content)
        else:
            sentences = content.split('. ')
            sentences = [s.strip() + '.' for s in sentences if s.strip()]
            
        if len(sentences) <= 3:
            return [self._create_chunk(content, filename, "Topic 1",
                                     metadata={'chunk_method': 'topic', 'chunk_index': 0})]
        
        try:
            # Simple topic detection using keyword overlap
            from collections import Counter
            import re
            
            # Extract keywords from each sentence
            sentence_keywords = []
            for sentence in sentences:
                # Simple keyword extraction (could be enhanced with NLP)
                words = re.findall(r'\b[a-zA-Z]{3,}\b', sentence.lower())
                # Filter common words (basic stopwords)
                stopwords = {'the', 'and', 'for', 'are', 'but', 'not', 'you', 'all', 'can', 'had', 'her', 'was', 'one', 'our', 'out', 'day', 'get', 'has', 'him', 'his', 'how', 'its', 'may', 'new', 'now', 'old', 'see', 'two', 'who', 'boy', 'did', 'man', 'way', 'she', 'use', 'her', 'many', 'oil', 'sit', 'set', 'run', 'eat', 'far', 'sea', 'eye', 'ask', 'own', 'say', 'too', 'any', 'try', 'us', 'an', 'as', 'at', 'be', 'he', 'if', 'in', 'is', 'it', 'my', 'of', 'on', 'or', 'to', 'up', 'we', 'go', 'no', 'so', 'am', 'by', 'do', 'me'}
                keywords = [w for w in words if w not in stopwords and len(w) > 3]
                sentence_keywords.append(set(keywords))
            
            # Find topic boundaries based on keyword overlap
            chunks = []
            current_chunk = [sentences[0]]
            current_keywords = sentence_keywords[0]
            
            for i in range(1, len(sentences)):
                # Calculate keyword overlap with current chunk
                overlap = len(current_keywords.intersection(sentence_keywords[i]))
                total_keywords = len(current_keywords.union(sentence_keywords[i]))
                
                if total_keywords > 0:
                    similarity = overlap / total_keywords
                else:
                    similarity = 0
                
                # If similarity is low, start new chunk
                if similarity < self.topic_threshold and len(current_chunk) >= 2:
                    chunk_content = ' '.join(current_chunk)
                    chunks.append(self._create_chunk(
                        content=chunk_content,
                        filename=filename,
                        section=f"Topic {len(chunks)+1}",
                        metadata={
                            'chunk_method': 'topic',
                            'chunk_index': len(chunks),
                            'topic_keywords': list(current_keywords)[:10],  # Top keywords
                            'sentence_count': len(current_chunk),
                            'topic_threshold': self.topic_threshold
                        }
                    ))
                    current_chunk = [sentences[i]]
                    current_keywords = sentence_keywords[i]
                else:
                    current_chunk.append(sentences[i])
                    current_keywords = current_keywords.union(sentence_keywords[i])
            
            # Add final chunk
            if current_chunk:
                chunk_content = ' '.join(current_chunk)
                chunks.append(self._create_chunk(
                    content=chunk_content,
                    filename=filename,
                    section=f"Topic {len(chunks)+1}",
                    metadata={
                        'chunk_method': 'topic',
                        'chunk_index': len(chunks),
                        'topic_keywords': list(current_keywords)[:10],
                        'sentence_count': len(current_chunk),
                        'topic_threshold': self.topic_threshold
                    }
                ))
            
            return chunks if chunks else [self._create_chunk(content, filename, "Topic 1",
                                                           metadata={'chunk_method': 'topic', 'chunk_index': 0})]
            
        except Exception:
            # Fallback to sentence-based chunking
            return self._chunk_by_sentences(content, filename, file_type)

    def _chunk_by_qa_optimization(self, content: str, filename: str, file_type: str) -> List[Dict[str, Any]]:
        """Create chunks optimized for question-answering"""
        if isinstance(content, list):
            content = '\n'.join(content)
        
        if sent_tokenize:
            sentences = sent_tokenize(content)
        else:
            sentences = content.split('. ')
            sentences = [s.strip() + '.' for s in sentences if s.strip()]
        
        # Patterns that indicate Q&A structure
        question_patterns = [
            r'\?',  # Questions
            r'^(what|how|why|when|where|who|which|can|does|is|are|will|would|should)',
            r'(step|steps|process|procedure|method|way to)',
            r'(example|examples|instance|case)',
            r'(definition|meaning|refers to|means)',
        ]
        
        chunks = []
        current_chunk = []
        current_context = []
        
        for i, sentence in enumerate(sentences):
            sentence_lower = sentence.lower().strip()
            
            # Check if this sentence contains Q&A indicators
            is_qa_relevant = any(re.search(pattern, sentence_lower) for pattern in question_patterns)
            
            if is_qa_relevant or len(current_chunk) == 0:
                current_chunk.append(sentence)
                # Add surrounding context (previous and next sentences)
                if i > 0 and sentences[i-1] not in current_chunk and sentences[i-1] not in current_context:
                    current_context.append(sentences[i-1])
                if i < len(sentences) - 1 and sentences[i+1] not in current_chunk and sentences[i+1] not in current_context:
                    current_context.append(sentences[i+1])
            else:
                current_chunk.append(sentence)
            
            # Create chunk when we have enough content or reach a natural break
            if (len(current_chunk) >= 3 and 
                (i == len(sentences) - 1 or  # Last sentence
                 sentence.endswith('.') and len(current_chunk) >= 5)):  # Natural break
                
                # Combine chunk with context
                full_content = current_context + current_chunk
                chunk_content = ' '.join(full_content)
                
                chunks.append(self._create_chunk(
                    content=chunk_content,
                    filename=filename,
                    section=f"QA Section {len(chunks)+1}",
                    metadata={
                        'chunk_method': 'qa_optimized',
                        'chunk_index': len(chunks),
                        'has_question': any('?' in s for s in current_chunk),
                        'has_process': any(re.search(r'(step|process|method)', s.lower()) for s in current_chunk),
                        'sentence_count': len(full_content)
                    }
                ))
                
                current_chunk = []
                current_context = []
        
        # Handle remaining content
        if current_chunk:
            chunk_content = ' '.join(current_context + current_chunk)
            chunks.append(self._create_chunk(
                content=chunk_content,
                filename=filename,
                section=f"QA Section {len(chunks)+1}",
                metadata={
                    'chunk_method': 'qa_optimized',
                    'chunk_index': len(chunks),
                    'sentence_count': len(current_context + current_chunk)
                }
            ))
        
        return chunks if chunks else [self._create_chunk(content, filename, "QA Section 1",
                                                       metadata={'chunk_method': 'qa_optimized', 'chunk_index': 0})]
    
    def _chunk_from_json(self, content: str, filename: str, file_type: str) -> List[Dict[str, Any]]:
        """
        Create chunks from pre-processed JSON content
        
        This strategy expects content to be a JSON string with the following structure:
        {
            "chunks": [
                {
                    "chunk_id": "unique_id",
                    "type": "content|toc",
                    "content": "text content",
                    "metadata": {
                        "url": "https://...",
                        "section_number": 1,
                        "related_toc": "toc_id",
                        ...
                    }
                },
                ...
            ]
        }
        
        Args:
            content: JSON string containing pre-chunked content
            filename: Name of the source file
            file_type: Should be 'json'
            
        Returns:
            List of chunk dictionaries formatted for the search index
        """
        try:
            # Parse JSON content
            data = json.loads(content)
            
            if not isinstance(data, dict) or 'chunks' not in data:
                logger.error(f"Invalid JSON structure in {filename}: expected 'chunks' key")
                # Fallback to treating it as plain text
                return self._chunk_by_sentences(content, filename, file_type)
            
            chunks = []
            for idx, json_chunk in enumerate(data['chunks']):
                if not isinstance(json_chunk, dict) or 'content' not in json_chunk:
                    logger.warning(f"Skipping invalid chunk {idx} in {filename}")
                    continue
                
                # Extract metadata from JSON chunk
                json_metadata = json_chunk.get('metadata', {})
                chunk_type = json_chunk.get('type', 'content')
                
                # Build chunk metadata (excluding tags which go at top level)
                metadata = {
                    'chunk_method': 'json',
                    'chunk_index': idx,
                    'chunk_type': chunk_type,
                    'original_chunk_id': json_chunk.get('chunk_id', f'chunk_{idx}')
                }
                
                # Extract tags before merging metadata
                tags = json_metadata.get('tags', [])
                
                # Merge JSON metadata (this includes all fields including tags)
                # We'll keep tags in metadata for backward compatibility but also set at top level
                metadata.update(json_metadata)
                
                # Determine section name
                if chunk_type == 'toc':
                    section = f"TOC: {json_chunk.get('content', '')[:50]}"
                else:
                    section = json_metadata.get('section', f"Section {json_metadata.get('section_number', idx + 1)}")
                
                # Create chunk with proper structure
                chunk = self._create_chunk(
                    content=json_chunk['content'],
                    filename=filename,
                    section=section,
                    metadata=metadata
                )
                
                # Set tags at the top level for proper tag filtering
                if tags:
                    chunk['tags'] = tags
                elif chunk_type == 'toc':
                    # For TOC entries, add special tags if none provided
                    chunk['tags'] = ['toc', 'navigation']
                
                chunks.append(chunk)
            
            if not chunks:
                logger.warning(f"No valid chunks found in JSON file {filename}")
                return self._chunk_by_sentences(str(data), filename, file_type)
            
            logger.info(f"Created {len(chunks)} chunks from JSON file {filename}")
            return chunks
            
        except json.JSONDecodeError as e:
            logger.error(f"Failed to parse JSON in {filename}: {e}")
            # Fallback to sentence chunking
            return self._chunk_by_sentences(content, filename, file_type)
        except Exception as e:
            logger.error(f"Unexpected error processing JSON chunks in {filename}: {e}")
            return self._chunk_by_sentences(content, filename, file_type) 